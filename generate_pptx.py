"""
Generate an 80-slide PowerPoint curriculum deck — one idea per slide, no cramming.
Run: python3 generate_pptx.py
Output: curriculum_advaith.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────
DARK       = RGBColor(0x1a, 0x1a, 0x2e)
MID        = RGBColor(0x16, 0x21, 0x3e)
CARD       = RGBColor(0x0f, 0x1a, 0x2e)
RED        = RGBColor(0xe9, 0x45, 0x60)
TEAL       = RGBColor(0x00, 0xd4, 0xaa)
GOLD       = RGBColor(0xff, 0xd7, 0x00)
PURPLE     = RGBColor(0x7c, 0x3a, 0xed)
BLUE       = RGBColor(0x08, 0x91, 0xb2)
GREEN      = RGBColor(0x05, 0x96, 0x69)
AMBER      = RGBColor(0xd9, 0x77, 0x06)
SLATE      = RGBColor(0x6b, 0x72, 0x80)
INSTA      = RGBColor(0xe1, 0x30, 0x6c)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LGREY      = RGBColor(0xcc, 0xcc, 0xdd)
MGREY      = RGBColor(0x88, 0x88, 0x99)

# Session colour map
C = {
    "intro":    TEAL,
    "s1":       RED,
    "s11":      RED,
    "s12":      RED,
    "s2":       PURPLE,
    "s3":       PURPLE,
    "s4":       BLUE,
    "s5":       BLUE,
    "s6":       GREEN,
    "s7":       AMBER,
    "s8":       GOLD,
    "s9":       SLATE,
    "s10":      INSTA,
}

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════════════
# PRIMITIVES
# ════════════════════════════════════════════════════════════════

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, col):
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = col

def rect(s, x, y, w, h, col, line=False):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    if not line:
        sh.line.fill.background()
    return sh

def oval(s, x, y, w, h, col):
    sh = s.shapes.add_shape(9, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, sz=18, bold=False, col=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    b = s.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col
    return b, tf

def add_p(tf, text, sz=14, bold=False, col=WHITE,
          align=PP_ALIGN.LEFT, italic=False, space=3):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col

def hline(s, y, col=RED, w=Inches(12.5), x=Inches(0.4)):
    l = s.shapes.add_shape(1, x, y, w, Pt(2))
    l.fill.solid()
    l.fill.fore_color.rgb = col
    l.line.fill.background()

def pill(s, label, x, y, col, sz=9):
    rect(s, x, y, Inches(0.06), Inches(0.32), col)
    txt(s, label, x + Inches(0.12), y, Inches(2.5), Inches(0.32),
        sz=sz, bold=True, col=col)

def badge(s, text, x, y, w, col, text_col=WHITE):
    rect(s, x, y, w, Inches(0.34), col)
    txt(s, text, x + Inches(0.08), y + Inches(0.03), w - Inches(0.12), Inches(0.3),
        sz=9, bold=True, col=text_col, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE TEMPLATES
# ════════════════════════════════════════════════════════════════

def make_cover():
    s = slide()
    bg(s, DARK)
    # Left accent bar
    rect(s, 0, 0, Inches(0.2), H, RED)
    # Decorative circles
    oval(s, Inches(10.5), Inches(-0.5), Inches(3.5), Inches(3.5), RGBColor(0x25,0x0a,0x12))
    oval(s, Inches(11.5), Inches(4.5), Inches(2.5), Inches(2.5), RGBColor(0x0a,0x1a,0x25))
    # Badge
    rect(s, Inches(0.5), Inches(1.0), Inches(3.0), Inches(0.36), RED)
    txt(s, "NITW M&C PREPARATION  ·  2026",
        Inches(0.5), Inches(1.0), Inches(3.0), Inches(0.36),
        sz=8.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    # Title
    txt(s, "Advaith's", Inches(0.5), Inches(1.55), Inches(9), Inches(0.6),
        sz=26, col=LGREY)
    txt(s, "2-Month Curriculum",
        Inches(0.5), Inches(2.1), Inches(10), Inches(1.1),
        sz=52, bold=True, col=WHITE)
    # Tagline
    txt(s, "Computing  ·  Math  ·  Python  ·  AI  ·  Quant Finance  ·  Software Engineering",
        Inches(0.5), Inches(3.3), Inches(11), Inches(0.45),
        sz=14, col=TEAL, italic=True)
    hline(s, Inches(3.9), TEAL, Inches(5))
    # Meta
    pairs = [("Student","Advaith"),("Teacher","Sree"),("Duration","2 Months"),("Sessions","11 + Capstone")]
    for i,(k,v) in enumerate(pairs):
        x = Inches(0.5 + i*2.0)
        txt(s, k.upper(), x, Inches(4.1), Inches(1.8), Inches(0.28), sz=8, col=MGREY, bold=True)
        txt(s, v, x, Inches(4.35), Inches(1.8), Inches(0.42), sz=17, col=WHITE, bold=True)
    # Sports emojis
    txt(s, "⚽  🏀  📊  🤖  📈  🏗️",
        Inches(0.5), Inches(5.45), Inches(9), Inches(0.75), sz=28)
    txt(s, "Taught by Sree  ·  Built for NITW M&C  ·  Starting from zero",
        Inches(0.5), Inches(6.65), Inches(9), Inches(0.45),
        sz=10, col=MGREY, italic=True)

def make_chapter_break(num, emoji, title, subtitle, hook, col):
    """Full-bleed session divider — one idea only: we're entering a new session."""
    s = slide()
    bg(s, DARK)
    # Big coloured strip left third
    rect(s, 0, 0, Inches(5.2), H, col)
    # Session number — huge (lighten the colour slightly for contrast)
    txt(s, num, Inches(0.25), Inches(0.9), Inches(4.7), Inches(3.5),
        sz=140, bold=True, col=RGBColor(
            min(col[0]+30, 255), min(col[1]+20, 255), min(col[2]+20, 255)))
    # Emoji
    txt(s, emoji, Inches(0.3), Inches(5.1), Inches(1.5), Inches(1.0), sz=42)
    # Title on right
    txt(s, "SESSION " + num, Inches(5.5), Inches(1.0), Inches(7.5), Inches(0.4),
        sz=10, bold=True, col=col)
    txt(s, title, Inches(5.5), Inches(1.45), Inches(7.5), Inches(1.3),
        sz=34, bold=True, col=WHITE)
    txt(s, subtitle, Inches(5.5), Inches(2.8), Inches(7.5), Inches(0.5),
        sz=15, col=LGREY, italic=True)
    hline(s, Inches(3.5), col, Inches(7.5), Inches(5.5))
    # Hook quote
    txt(s, hook, Inches(5.5), Inches(3.75), Inches(7.5), Inches(2.0),
        sz=14, col=WHITE, italic=True)

def make_big_quote(quote, who, col):
    """A slide that is just one powerful quote."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, Inches(2.8), W, Inches(0.06), col)
    txt(s, "“", Inches(0.6), Inches(0.5), Inches(2), Inches(1.8),
        sz=120, bold=True, col=col)
    txt(s, quote, Inches(1.2), Inches(1.4), Inches(11.5), Inches(3.5),
        sz=22, col=WHITE, italic=True)
    txt(s, f"— {who}", Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
        sz=13, col=col, italic=True)

def make_content(title, bullets, col, note=None, sub=None):
    """Standard content slide: title + up to 6 bullets."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, col)
    if sub:
        txt(s, sub, Inches(0.35), Inches(0.25), Inches(12), Inches(0.35),
            sz=9.5, col=col, bold=True)
    txt(s, title, Inches(0.35), Inches(0.5 if sub else 0.3),
        Inches(12.5), Inches(0.85), sz=30, bold=True, col=WHITE)
    hline(s, Inches(1.42), col, Inches(12.8), Inches(0.35))
    # Bullets
    bx, tf = txt(s, "", Inches(0.5), Inches(1.65), Inches(12.0), Inches(5.2))
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = b
        r.font.size = Pt(18)
        r.font.color.rgb = WHITE
    # Optional note at bottom
    if note:
        rect(s, Inches(0.35), Inches(6.55), Inches(12.6), Inches(0.62), CARD)
        rect(s, Inches(0.35), Inches(6.55), Inches(0.08), Inches(0.62), col)
        txt(s, note, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.55),
            sz=11, col=LGREY, italic=True)

def make_two_col(title, left_head, left_items, right_head, right_items,
                 col, left_col=None, right_col=None, sub=None):
    """Two side-by-side columns."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, col)
    if sub:
        txt(s, sub, Inches(0.35), Inches(0.22), Inches(12), Inches(0.32), sz=9, col=col, bold=True)
    txt(s, title, Inches(0.35), Inches(0.45 if sub else 0.28),
        Inches(12.5), Inches(0.8), sz=30, bold=True, col=WHITE)
    hline(s, Inches(1.38), col, Inches(12.8), Inches(0.35))
    lc = left_col  or col
    rc = right_col or TEAL
    for idx, (head, items, accent) in enumerate(
            [(left_head, left_items, lc), (right_head, right_items, rc)]):
        x = Inches(0.38 + idx * 6.48)
        rect(s, x, Inches(1.58), Inches(6.1), Inches(5.55), MID)
        rect(s, x, Inches(1.58), Inches(0.09), Inches(5.55), accent)
        txt(s, head, x + Inches(0.2), Inches(1.66), Inches(5.7), Inches(0.4),
            sz=13, bold=True, col=accent)
        bx, tf = txt(s, "", x + Inches(0.18), Inches(2.15), Inches(5.7), Inches(4.7))
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(8)
            r = p.add_run()
            r.text = f"▸  {item}"
            r.font.size = Pt(14)
            r.font.color.rgb = WHITE

def make_lab(session_label, lab_title, steps, output_desc, col, sports_note=None):
    """Lab slide: what to build, step by step."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, GREEN)
    # Header
    badge(s, f"🔧  LAB  ·  {session_label}", Inches(0.35), Inches(0.22), Inches(4.0), GREEN)
    txt(s, lab_title, Inches(0.35), Inches(0.7), Inches(12.5), Inches(0.82),
        sz=28, bold=True, col=WHITE)
    hline(s, Inches(1.6), GREEN, Inches(12.8), Inches(0.35))
    # Steps (left)
    rect(s, Inches(0.35), Inches(1.78), Inches(7.3), Inches(4.65), MID)
    rect(s, Inches(0.35), Inches(1.78), Inches(0.09), Inches(4.65), GREEN)
    txt(s, "STEPS", Inches(0.55), Inches(1.85), Inches(7.0), Inches(0.3),
        sz=9, bold=True, col=GREEN)
    bx, tf = txt(s, "", Inches(0.52), Inches(2.2), Inches(7.0), Inches(4.1))
    tf.word_wrap = True
    first = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(7)
        r = p.add_run()
        r.text = f"{i+1}.  {step}"
        r.font.size = Pt(14)
        r.font.color.rgb = WHITE
    # Output (right)
    rect(s, Inches(7.85), Inches(1.78), Inches(5.1), Inches(2.9), RGBColor(0x05,0x20,0x14))
    rect(s, Inches(7.85), Inches(1.78), Inches(0.09), Inches(2.9), TEAL)
    txt(s, "✓  OUTPUT", Inches(8.05), Inches(1.85), Inches(4.8), Inches(0.3),
        sz=9, bold=True, col=TEAL)
    txt(s, output_desc, Inches(8.05), Inches(2.2), Inches(4.8), Inches(2.35),
        sz=13, col=WHITE)
    # Sports note
    if sports_note:
        rect(s, Inches(7.85), Inches(4.85), Inches(5.1), Inches(1.6), RGBColor(0x20,0x10,0x05))
        rect(s, Inches(7.85), Inches(4.85), Inches(0.09), Inches(1.6), RED)
        txt(s, "⚽🏀  SPORTS", Inches(8.05), Inches(4.92), Inches(4.8), Inches(0.3),
            sz=9, bold=True, col=RED)
        txt(s, sports_note, Inches(8.05), Inches(5.28), Inches(4.8), Inches(1.1),
            sz=12, col=WHITE)

def make_assignment(session_label, tasks, col):
    """Assignment slide: what Advaith does independently."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, AMBER)
    badge(s, f"📝  ASSIGNMENT  ·  {session_label}", Inches(0.35), Inches(0.22), Inches(5.0), AMBER)
    txt(s, "Independent Work", Inches(0.35), Inches(0.7), Inches(12.5), Inches(0.8),
        sz=30, bold=True, col=WHITE)
    txt(s, "Complete before the next session. Bring all answers.",
        Inches(0.35), Inches(1.48), Inches(12.0), Inches(0.38),
        sz=12, col=LGREY, italic=True)
    hline(s, Inches(1.93), AMBER, Inches(12.8), Inches(0.35))
    # Task cards — 2 per row, up to 4 tasks
    task_color = [col, TEAL, PURPLE, GREEN]
    positions = [
        (Inches(0.38), Inches(2.12)),
        (Inches(6.72), Inches(2.12)),
        (Inches(0.38), Inches(4.8)),
        (Inches(6.72), Inches(4.8)),
    ]
    for i, (task, pos) in enumerate(zip(tasks, positions)):
        x, y = pos
        w, h = Inches(6.0), Inches(2.45)
        tc = task_color[i % len(task_color)]
        rect(s, x, y, w, h, MID)
        rect(s, x, y, Inches(0.09), h, tc)
        rect(s, x, y, w, Inches(0.08), tc)
        txt(s, f"Task {i+1}", x + Inches(0.18), y + Inches(0.1),
            Inches(5.7), Inches(0.3), sz=9, bold=True, col=tc)
        txt(s, task, x + Inches(0.18), y + Inches(0.42), Inches(5.7), Inches(1.85),
            sz=12.5, col=WHITE)

def make_comparison(title, left_label, left_items, right_label, right_items,
                    col, left_col=RED, right_col=GREEN, sub=None):
    """A vs B / Before vs After comparison slide."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, col)
    if sub:
        txt(s, sub, Inches(0.35), Inches(0.22), Inches(12), Inches(0.3), sz=9, col=col, bold=True)
    txt(s, title, Inches(0.35), Inches(0.45 if sub else 0.28),
        Inches(12.5), Inches(0.8), sz=30, bold=True, col=WHITE)
    hline(s, Inches(1.38), col, Inches(12.8), Inches(0.35))
    for idx, (label, items, accent) in enumerate(
            [(left_label, left_items, left_col), (right_label, right_items, right_col)]):
        x = Inches(0.38 + idx * 6.48)
        rect(s, x, Inches(1.58), Inches(6.1), Inches(5.6), MID)
        rect(s, x, Inches(1.58), Inches(6.1), Inches(0.48), accent)
        txt(s, label, x + Inches(0.2), Inches(1.63), Inches(5.7), Inches(0.38),
            sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        bx, tf = txt(s, "", x + Inches(0.2), Inches(2.18), Inches(5.7), Inches(4.8))
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(10)
            r = p.add_run()
            r.text = f"{'✗' if accent==left_col else '✓'}  {item}"
            r.font.size = Pt(15)
            r.font.color.rgb = RGBColor(0xee,0x88,0x88) if accent==left_col else RGBColor(0x88,0xee,0xaa)

def make_code_slide(title, code_lines, explanation, col, sub=None):
    """A slide showing a short code snippet with explanation."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, col)
    if sub:
        txt(s, sub, Inches(0.35), Inches(0.22), Inches(12), Inches(0.3), sz=9, col=col, bold=True)
    txt(s, title, Inches(0.35), Inches(0.45 if sub else 0.28),
        Inches(12.5), Inches(0.78), sz=28, bold=True, col=WHITE)
    hline(s, Inches(1.32), col, Inches(12.8), Inches(0.35))
    # Code block
    rect(s, Inches(0.35), Inches(1.52), Inches(7.6), Inches(5.65), RGBColor(0x1e,0x1e,0x2e))
    rect(s, Inches(0.35), Inches(1.52), Inches(0.09), Inches(5.65), col)
    bx, tf = txt(s, "", Inches(0.52), Inches(1.65), Inches(7.3), Inches(5.4))
    tf.word_wrap = False
    first = True
    for line in code_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12.5)
        r.font.name = "Courier New"
        r.font.color.rgb = RGBColor(0xcd, 0xd6, 0xf4)
    # Explanation
    rect(s, Inches(8.12), Inches(1.52), Inches(4.88), Inches(5.65), MID)
    rect(s, Inches(8.12), Inches(1.52), Inches(0.09), Inches(5.65), TEAL)
    txt(s, "WHAT IT DOES", Inches(8.3), Inches(1.6), Inches(4.55), Inches(0.3),
        sz=9, bold=True, col=TEAL)
    txt(s, explanation, Inches(8.3), Inches(2.0), Inches(4.55), Inches(5.0),
        sz=13.5, col=WHITE)

def make_timeline(title, items, col, sub=None):
    """Vertical timeline slide."""
    s = slide()
    bg(s, DARK)
    rect(s, 0, 0, Inches(0.15), H, col)
    if sub:
        txt(s, sub, Inches(0.35), Inches(0.22), Inches(12), Inches(0.3), sz=9, col=col, bold=True)
    txt(s, title, Inches(0.35), Inches(0.45 if sub else 0.28),
        Inches(12.5), Inches(0.75), sz=30, bold=True, col=WHITE)
    hline(s, Inches(1.3), col, Inches(12.8), Inches(0.35))
    # Vertical line
    rect(s, Inches(1.75), Inches(1.5), Inches(0.06), Inches(5.7), MGREY)
    row_h = Inches(5.7) / max(len(items), 1)
    for i, (year, title_t, desc) in enumerate(items):
        y = Inches(1.5) + i * row_h + Inches(0.08)
        # Dot on timeline
        oval(s, Inches(1.6), y + Inches(0.08), Inches(0.32), Inches(0.32), col)
        # Year badge
        rect(s, Inches(0.35), y, Inches(1.15), Inches(0.36), col)
        txt(s, year, Inches(0.35), y, Inches(1.15), Inches(0.36),
            sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        # Title + desc
        txt(s, title_t, Inches(2.25), y, Inches(10.6), Inches(0.32),
            sz=14, bold=True, col=WHITE)
        txt(s, desc, Inches(2.25), y + Inches(0.3), Inches(10.6), Inches(0.3),
            sz=11, col=LGREY, italic=True)


# ════════════════════════════════════════════════════════════════
# ALL 80 SLIDES
# ════════════════════════════════════════════════════════════════

print("🏗️  Building 80-slide PowerPoint deck...\n")

# ── INTRO (slides 1–6) ──────────────────────────────────────────
print("  [Intro] Slides 1–6")

# 1. Cover
make_cover()

# 2. About This Curriculum
make_content("What This Curriculum Is", [
    "⏱️  Two months of structured learning — one dedicated session at a time",
    "👨‍🏫  Sree teaches every session. Advaith practices independently after each.",
    "🎯  Goal: arrive at NITW M&C with genuine understanding, not surface familiarity",
    "🧱  Builds from zero — no prior coding knowledge assumed",
    "⚽🏀  Basketball and soccer woven into every session as the primary examples",
    "🏆  Ends with a real, working capstone project Advaith presents himself",
], TEAL)

# 3. Advaith's Profile
make_two_col(
    "Who Is Advaith?",
    "📚 Starting Point",
    ["Just finished 12th PCM",
     "Loves math — scored well, genuinely enjoyed it",
     "Strong in calculus, vectors, probability, stats",
     "Never written a line of code",
     "Zero programming experience — blank slate"],
    "🎯 By the End",
    ["Can write and explain Python programs",
     "Understands how AI works at a mathematical level",
     "Has built a working AI-powered sports analytics tool",
     "Knows the full software engineering landscape",
     "Ready for NITW M&C from day one"],
    TEAL
)

# 4. The Teaching Philosophy
make_content("Every Session Has the Same Shape", [
    "🎬  HOOK — Sree opens with a story or demo that creates a question",
    "📖  EXPLAIN — Sree walks through the concept live, no slides-only lectures",
    "🔧  DEMO — Sree runs real code or a visualisation",
    "🤝  PRACTICE — Sree and Advaith build it together",
    "⚡  WOW — the moment everything clicks into place",
    "📝  HOMEWORK — Advaith continues independently before the next session",
], TEAL, note="Rule: no session starts with 'today we learn X'. Every session starts with something that makes Advaith want to understand X.")

# 5. Curriculum Map
s = slide()
bg(s, DARK)
rect(s, 0, 0, Inches(0.15), H, TEAL)
txt(s, "The Full Journey — 11 Sessions", Inches(0.35), Inches(0.28),
    Inches(12.5), Inches(0.78), sz=30, bold=True, col=WHITE)
hline(s, Inches(1.18), TEAL)
sessions_map = [
    ("1",   "🕰️", "History of Computing",          C["s1"]),
    ("1.1", "⚡", "Building Blocks — CPU to GPU",   C["s11"]),
    ("1.2", "🌐", "How the Web Works",              C["s12"]),
    ("2",   "🔢", "Maths for Computing",            C["s2"]),
    ("3",   "🌍", "Programming Languages",          C["s3"]),
    ("4",   "🐍", "Python — Math with Superpowers", C["s4"]),
    ("5",   "🧠", "Math & Stats for AI",            C["s5"]),
    ("6",   "🤖", "Agentic AI",                     C["s6"]),
    ("7",   "📈", "Quant Finance",                  C["s7"]),
    ("8",   "🏆", "Capstone Project",               C["s8"]),
    ("9",   "🏗️", "Software Engineering Landscape", C["s9"]),
    ("10",  "📸", "How Instagram Works",           C["s10"]),
]
row_h = Inches(0.52)
cols = 2
for i,(num,em,title,col) in enumerate(sessions_map):
    col_idx = i % cols
    row_idx = i // cols
    x = Inches(0.38 + col_idx*6.45)
    y = Inches(1.35) + row_idx * row_h
    rect(s, x, y+Inches(0.07), Inches(0.46), Inches(0.34), col)
    txt(s, num, x, y+Inches(0.07), Inches(0.46), Inches(0.34),
        sz=9, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    txt(s, f"{em}  {title}", x+Inches(0.54), y+Inches(0.08), Inches(5.7), Inches(0.34),
        sz=13, bold=True, col=WHITE)

# 6. 8-Week Roadmap
make_timeline("8-Week Learning Roadmap", [
    ("Wk 1",  "History + Hardware + Web",         "Sessions 1, 1.1, 1.2  ·  No coding yet — just understanding"),
    ("Wk 2",  "Maths + Languages",               "Sessions 2 & 3  ·  Caesar cipher, passing networks, Fibonacci"),
    ("Wk 3",  "Python — First Half",             "Session 4 (part 1)  ·  Waves, guessing game"),
    ("Wk 4",  "Python — Second Half",            "Session 4 (part 2)  ·  NBA charts, dice, law of large numbers"),
    ("Wk 5",  "Math & Stats for AI",             "Session 5  ·  Vectors, gradient descent, xG model"),
    ("Wk 6",  "Agentic AI",                      "Session 6  ·  Build the sports analyst agent"),
    ("Wk 7",  "Quant Finance + SE Landscape",    "Sessions 7 & 9  ·  Man United stock, 2010→2026 map"),
    ("Wk 8",  "Capstone + Instagram Deep Dive", "Sessions 8 & 10  ·  Build day + how Instagram really works"),
], TEAL)

# ── SESSION 1: HISTORY (slides 7–13) ────────────────────────────
print("  [Session 1] Slides 7–13")

# 7. Divider
make_chapter_break("1", "🕰️", "History of Computing",
    "The Greatest Heist in History",
    '"The first programmer was a woman in 1843.\nThe machine she programmed didn\'t exist yet."',
    C["s1"])

# 8. Ada Lovelace
make_content("Ada Lovelace — The First Programmer", [
    "👩‍💻  Charles Babbage designed the Difference Engine in the 1830s — a mechanical calculator",
    "📝  Ada Lovelace wrote an algorithm for it in 1843 — for a machine that wasn't finished",
    "🧮  Her algorithm computed Bernoulli numbers — a real mathematical sequence",
    "💡  She saw that the machine could do more than arithmetic — it could manipulate symbols",
    "🏆  She is recognised as the world's first computer programmer — 100 years before computers existed",
], C["s1"], note="Sree: 'She wrote an algorithm for a machine that didn't fully exist yet. What does that tell you about the relationship between ideas and technology?'",
sub="Session 1 — History of Computing")

# 9. Alan Turing
make_content("Alan Turing — The Greatest Codebreaker", [
    "🔐  World War 2: Nazi Germany used the Enigma machine to encrypt all military communications",
    "🖥️  Turing built the Bombe machine — cracked Enigma, turning the tide of the war",
    "📊  Historians estimate he shortened the war by 2 years and saved 14 million lives",
    "🤔  He asked: 'Can machines think?' — invented the theoretical Turing Machine (basis of all computing)",
    "💔  After the war, the British government prosecuted him for being gay. He died at 41.",
    "🌍  The Imitation Game (2014 film) tells his story. Advaith's homework: watch it.",
], C["s1"], sub="Session 1 — History of Computing")

# 10. Machines Arrive
make_timeline("The Machines Arrive", [
    ("1945", "ENIAC — First General-Purpose Computer",  "Filled a room · 30 tonnes · 18,000 vacuum tubes · 5,000 additions/second"),
    ("1947", "Transistor Invented at Bell Labs",         "Silicon switch replaces vacuum tubes · Smaller, faster, reliable"),
    ("1958", "Integrated Circuit",                       "Thousands of transistors on one chip · Moore's Law begins"),
    ("1965", "Moore's Law",                              "Gordon Moore: transistors double every ~2 years · Held for 60 years"),
    ("1971", "Intel 4004 — First Microprocessor",        "CPU on a single chip · Paved the way for personal computers"),
], C["s1"], sub="Session 1 — History of Computing")

# 11. Personal Computer + Web
make_timeline("The World Shrinks", [
    ("1977", "Apple II — Personal Computer",      "Computers leave corporations · Enter homes · Steve Jobs, Steve Wozniak"),
    ("1984", "Apple Macintosh",                   "Mouse + windows + icons · Made computers human · 1984 ad is legendary"),
    ("1989", "World Wide Web",                    "Tim Berners-Lee invents HTML + HTTP + URLs · Gave it away for free"),
    ("1998", "Google Search",                     "Algorithm ranks web pages by importance · PageRank = graph centrality"),
    ("2007", "iPhone",                            "'An iPod, a phone, and an internet communicator' · The computer becomes personal"),
], C["s1"], sub="Session 1 — History of Computing")

# 12. AI + Sports Analytics Revolution
make_two_col(
    "Two Revolutions Converge",
    "🤖 The AI Revolution",
    ["2012: AlexNet wins ImageNet — deep learning works",
     "2017: 'Attention Is All You Need' — Transformers invented",
     "2020: GPT-3 — scale changes everything",
     "2022: ChatGPT — 1 million users in 5 days",
     "2025: AI agents take autonomous actions"],
    "⚽🏀 The Sports Analytics Revolution",
    ["2002: Moneyball — Oakland A's use stats to win on half the payroll",
     "2012: Daryl Morey proves mid-range shots are the worst shot in basketball",
     "2015: Every NBA team installs player-tracking cameras (Second Spectrum)",
     "2018: StatsBomb xG model changes how football is analysed",
     "Today: Every top club has a data science team"],
    C["s1"], left_col=C["s1"], right_col=TEAL,
    sub="Session 1 — History of Computing"
)

# 13. Lab + Assignment
make_lab("Session 1", "The Computing & Sports Analytics Timeline",
    ["Research 5 computing milestones (Sree provides starting points)",
     "Research 5 sports analytics milestones (Advaith finds these)",
     "For each: write what happened + why it mattered (1 sentence each)",
     "Present the sports milestones to Sree — explain each one"],
    "A 10-row table: half computing, half sports. Advaith presents his 5 sports milestones aloud.",
    C["s1"],
    sports_note="Moneyball · NBA 3-point revolution · xG in football · GPS vests in training")

# ── SESSION 1.1: BUILDING BLOCKS (slides 14–19) ─────────────────
print("  [Session 1.1] Slides 14–19")

# 14. Divider
make_chapter_break("1.1", "⚡", "Building Blocks",
    "What's Actually Inside the Box",
    '"There are 19 billion transistors in an Apple M3 chip.\nEach one is 3 nanometres wide. A hair is 80,000nm."',
    C["s11"])

# 15. The Transistor — Everything Starts Here
make_content("The Transistor — It All Starts Here", [
    "💡  A transistor is a switch made of silicon — it has two states: ON (1) and OFF (0)",
    "🔋  Stack millions of switches together → represent any number in binary",
    "📦  Stack binary numbers together → represent any instruction",
    "📱  Stack billions of instructions → run Instagram",
    "🦠  A Covid virus is ~100 nanometres. Apple M3 transistors are 3nm. Smaller by 33×.",
    "📊  Your phone has more computing power than NASA had for the entire Apollo programme",
], C["s11"], note="Sree: show the electron microscope die shot of an Apple M3 chip. Let it land silently for 5 seconds.",
sub="Session 1.1 — Building Blocks")

# 16. Logic Gates
make_code_slide("Three Gates That Run the World",
    ["def AND(a, b):  return a and b",
     "def OR(a, b):   return a or b",
     "def NOT(a):     return not a",
     "",
     "# Test them",
     "AND(1, 1)  →  True",
     "AND(1, 0)  →  False",
     "OR(1, 0)   →  True",
     "OR(0, 0)   →  False",
     "NOT(1)     →  False",
     "NOT(0)     →  True",
     ],
    "Every decision a computer makes — every if/else in every app ever written — reduces to these three operations.\n\nAND, OR, NOT are Boolean algebra.\n\nAdvaith already knows this from 12th-grade set theory:\nAND = ∩\nOR = ∪\nNOT = complement",
    C["s11"], sub="Session 1.1 — Building Blocks")

# 17. CPU Architecture
make_content("How a CPU Works — The Fetch-Decode-Execute Cycle", [
    "1️⃣  FETCH — Get the next instruction from memory",
    "2️⃣  DECODE — Figure out what the instruction means",
    "3️⃣  EXECUTE — Do it: add numbers, store a value, jump to an address",
    "🔁  Repeat billions of times per second (clock speed = GHz = billion cycles/sec)",
    "🧮  ALU (Arithmetic Logic Unit) — the part that does the actual math",
    "💾  Memory hierarchy: Registers → Cache → RAM → SSD (fastest to slowest, smallest to largest)",
], C["s11"], note="Analogy: a chef reading a recipe (fetch), understanding a step (decode), chopping the onion (execute) — 3 billion times a second.",
sub="Session 1.1 — Building Blocks")

# 18. CPU vs GPU
make_comparison(
    "CPU vs GPU — Why AI Needs the Team",
    "🏀 CPU — The Superstar",
    ["8–16 powerful cores",
     "Can handle any complex task",
     "Switches context instantly",
     "Best for sequential tasks",
     "Like LeBron James — brilliant at everything",
     "One thing at a time, incredibly well"],
    "⚽ GPU — The Full Team",
    ["Thousands of small parallel cores",
     "NVIDIA H100: 16,896 CUDA cores",
     "All cores run simultaneously",
     "Best for matrix multiplication at scale",
     "Like 11 footballers executing a set play",
     "GPT-4 trained on 25,000+ GPUs for months"],
    C["s11"], sub="Session 1.1 — Building Blocks"
)

# 19. Lab + Assignment
make_lab("Session 1.1", "Logic Gates in Python",
    ["Open labs/logic_gates.py and run it as-is",
     "Predict each output BEFORE running — write your guess first",
     "Part 2: print truth tables for AND and OR (all 4 combinations)",
     "Part 3: build XOR from AND + OR + NOT",
     "Challenge: add NAND and NOR gates"],
    "Working logic gate functions + printed truth tables.\n\nChallenge: NAND and NOR with their own truth tables.",
    C["s11"],
    sports_note="CPU = LeBron · GPU = Barcelona pressing — same concept, different tools")

# ── SESSION 1.2: WEB (slides 20–25) ─────────────────────────────
print("  [Session 1.2] Slides 20–25")

# 20. Divider
make_chapter_break("1.2", "🌐", "How the Web Works",
    "50 Things Happen in Half a Second",
    '"What actually happens when you type youtube.com and press Enter?\nThere are 8 distinct steps — none of them magic."',
    C["s12"])

# 21. DNS + TCP
make_content("Step 1 & 2 — Find the Server and Connect", [
    "🗺️  DNS (Domain Name System) — converts 'youtube.com' → '142.250.80.46'",
    "📖  DNS is the internet's phone book — human names mapped to numeric addresses",
    "🔢  Every device on the internet has an IP address — its true identity",
    "📞  TCP — your browser and the server 'pick up the phone' — agree to talk",
    "🔒  HTTPS = encrypted TCP — both sides agree on a secret code before speaking",
    "🔐  The padlock in your browser means nobody can eavesdrop on this conversation",
], C["s12"], note="Ask Advaith: 'What's your phone number in this analogy?' Answer: your device's IP address.",
sub="Session 1.2 — How the Web Works")

# 22. HTTP Request/Response
make_code_slide("Step 3 & 4 — Asking and Answering",
    ["# The REQUEST your browser sends:",
     "GET /home HTTP/1.1",
     "Host: youtube.com",
     "User-Agent: Chrome/120",
     "Accept: text/html",
     "",
     "# The RESPONSE the server sends back:",
     "HTTP/1.1 200 OK",
     "Content-Type: text/html",
     "",
     "<!DOCTYPE html>",
     "<html>...the page...</html>",
    ],
    "Status codes:\n200 = OK, here you go\n301 = Moved, go here instead\n404 = Not found\n403 = Forbidden\n500 = Server crashed\n\nMethods:\nGET = fetch something\nPOST = send data\nDELETE = remove something\nPUT = update something",
    C["s12"], sub="Session 1.2 — How the Web Works")

# 23. HTML + CSS + JS
make_two_col(
    "Step 5 — Building the Page: Three Languages",
    "🦴 HTML — The Skeleton",
    ["Defines structure and content",
     "<h1> headings, <p> paragraphs",
     "<table> for data tables",
     "<a> for links, <img> for images",
     "The browser reads this and knows what exists on the page"],
    "👗 CSS — The Style",
    ["Controls how everything looks",
     "Colours, fonts, sizes, spacing",
     "background-color: #1a1a2e",
     "The dark theme on Advaith's sports page",
     "JavaScript — the muscles: makes things interactive and dynamic"],
    C["s12"], sub="Session 1.2 — How the Web Works"
)

# 24. APIs + JSON
make_code_slide("APIs — Apps Talking to Apps",
    ["# What an API response looks like (JSON):",
     "{",
     '  "player": "Vinícius Jr",',
     '  "goals": 24,',
     '  "assists": 11,',
     '  "club": "Real Madrid",',
     '  "xG": 18.3',
     "}",
     "",
     "# Your sports_analyst.py does this:",
     "client.messages.create(...)  # API call",
     "# Returns JSON → Python dict",
    ],
    "An API is a server that returns data (JSON) instead of a webpage.\n\nWhen Instagram loads your feed, it makes dozens of API calls:\n• one for posts\n• one for stories\n• one for suggested accounts\n• one for ads\n\nJavaScript assembles the JSON responses into what you see.\n\nThe UI is just a rendering of API data.",
    C["s12"], sub="Session 1.2 — How the Web Works")

# 25. Lab + Assignment
make_lab("Session 1.2", "First Webpage + Spy on the Internet",
    ["Open labs/my_sports_page.html in VS Code",
     "Replace all players with Advaith's actual favourites",
     "Change colours to a favourite team's hex code",
     "Add a third section (match results or anything sports)",
     "Open Chrome DevTools → Network tab → refresh Instagram",
     "Find one API/XHR call and look at the JSON response"],
    "A personalised sports webpage AND\na screenshot of DevTools showing live API calls.",
    C["s12"],
    sports_note="Sports page has soccer + basketball tables · DevTools shows how ESPN or the NBA app fetches live scores")

# ── SESSION 2: MATHS (slides 26–33) ─────────────────────────────
print("  [Session 2] Slides 26–33")

# 26. Divider
make_chapter_break("2", "🔢", "Maths for Computing",
    "Math Is the Operating System of Reality",
    '"You\'ve been doing computing math your whole life.\nYou just didn\'t know the names."',
    C["s2"])

# 27. Binary + Hex
make_content("Number Systems — Why Computers Count Differently", [
    "🔟  Humans use decimal (base-10) — because we have 10 fingers",
    "⚡  Computers use binary (base-2) — because transistors have 2 states: on/off",
    "📐  Each binary digit (bit) is a power of 2: 1, 2, 4, 8, 16, 32, 64...",
    "🔢  Decimal 13 = 8+4+1 = 1101 in binary",
    "🎨  Hexadecimal (base-16) — used in colours: #e94560 = red 233, green 69, blue 96",
    "💻  Every photo, video, text message is binary under the hood — just 1s and 0s",
], C["s2"], note="Hex shortcut: every pair of hex digits = one byte = 0-255. #ff0000 = pure red (255, 0, 0).",
sub="Session 2 — Maths for Computing")

# 28. Boolean Algebra
make_comparison(
    "Boolean Algebra = Set Theory You Already Know",
    "📚 12th Grade Set Theory",
    ["A ∩ B  — intersection",
     "A ∪ B  — union",
     "A'      — complement",
     "De Morgan: (A ∩ B)' = A' ∪ B'",
     "De Morgan: (A ∪ B)' = A' ∩ B'",
     "Venn diagrams"],
    "💻 Boolean Algebra in Computing",
    ["A AND B  — both must be true",
     "A OR B   — at least one must be true",
     "NOT A    — flip the value",
     "NOT(AND(A,B)) = OR(NOT(A), NOT(B))",
     "NOT(OR(A,B))  = AND(NOT(A), NOT(B))",
     "Logic gate circuits"],
    C["s2"],
    left_col=PURPLE, right_col=BLUE,
    sub="Session 2 — Maths for Computing"
)

# 29. Big-O Notation
make_content("Big-O — How Algorithms Scale", [
    "📏  Big-O describes how runtime grows as input size (n) grows",
    "🟢  O(1) — Constant: array lookup by index. Always instant. n doesn't matter.",
    "🟡  O(log n) — Logarithmic: binary search. 1 billion records → 30 guesses.",
    "🟠  O(n) — Linear: scan a list once. Double the list → double the time.",
    "🔴  O(n²) — Quadratic: compare every pair. 10× input → 100× time. Dangerous.",
    "💥  O(2ⁿ) — Exponential: brute-force chess. Grows impossibly fast.",
], C["s2"], note="Binary search is O(log n): 1,000 items → 10 guesses · 1,000,000 items → 20 guesses · 1 billion items → 30 guesses. Doubling the input costs ONE extra guess.",
sub="Session 2 — Maths for Computing")

# 30. Graph Theory
make_content("Graph Theory — The Math of Networks", [
    "🕸️  A graph = nodes (vertices) connected by edges",
    "🔢  Degree = how many connections a node has",
    "📍  Path = a route from one node to another",
    "⚡  Shortest path = fastest route (Dijkstra's algorithm, 1956)",
    "🏆  Centrality = how 'important' a node is — how many paths pass through it",
    "📱  Google Maps, social networks, the web, passing networks in football — all graphs",
], C["s2"], note="Sree: 'Google's PageRank algorithm that made them $100B is betweenness centrality on the web's link graph. Invented in 1956, used for billions of searches today.'",
sub="Session 2 — Maths for Computing")

# 31. Dijkstra's + Football Passing Networks
make_two_col(
    "Graph Theory in the Real World",
    "🗺️ Dijkstra's Algorithm",
    ["Written in 20 minutes in an Amsterdam café, 1956",
     "Edsger Dijkstra had no pen and paper — did it in his head",
     "Waited 3 years to publish — thought it was trivial",
     "Now runs in Google Maps millions of times per second",
     "Every GPS navigation uses this algorithm",
     "One of the most used algorithms in existence"],
    "⚽ Football Passing Networks",
    ["Players = nodes, passes = edges, pass count = weight",
     "Centrality reveals the team's key playmaker",
     "Xavi: highest centrality in Barcelona's tiki-taka system",
     "Remove high-centrality player → team flow collapses",
     "Same math as Google PageRank and social network analysis",
     "Clubs pay for this analysis in real scouting"],
    C["s2"], left_col=C["s2"], right_col=TEAL,
    sub="Session 2 — Maths for Computing"
)

# 32. Lab: Caesar Cipher
make_lab("Session 2 — Lab A", "Secret Message Machine — Caesar Cipher",
    ["Open labs/caesar_cipher.py — run it and study the output",
     "Understand how ord() and chr() convert letters to numbers",
     "Encrypt 'MESSI IS GOAT' with shift=3 — predict output first",
     "Decrypt the result back — verify you get the original",
     "Send Sree an encrypted message — see if he can decode it"],
    "encrypt() and decrypt() functions working.\nSend Sree a real secret message.",
    C["s2"],
    sports_note="Encrypt 'ADVAITH WILL ACE NITW' and send it to Sree")

# 33. Lab: Passing Network + Assignment
make_lab("Session 2 — Lab B", "Barcelona Passing Network",
    ["Open labs/passing_network.py — install networkx first",
     "Run Part 1: build the graph from the passing data",
     "Run Part 2: calculate betweenness centrality for all players",
     "Run Part 3: visualise the network — save the chart",
     "Remove Xavi from the data and re-run. Who becomes central now?"],
    "A visualised passing network chart saved as an image.\nCentrality ranking printed. Discussion: what does it mean tactically?",
    C["s2"],
    sports_note="Who is Barcelona's most important player mathematically? Does it match your intuition?")

# ── SESSION 3: LANGUAGES (slides 34–38) ─────────────────────────
print("  [Session 3] Slides 34–38")

# 34. Divider
make_chapter_break("3", "🌍", "Programming Languages",
    "The Family Tree of Code",
    '"Every programming language was invented to solve someone\'s frustration.\nHere\'s the 70-year drama."',
    C["s3"])

# 35. The Family Tree
make_timeline("The Language Family Tree", [
    ("1954", "FORTRAN — IBM",          "Scientists wanted to write math, not binary · 'Formula Translation'"),
    ("1972", "C — Bell Labs",          "Dennis Ritchie needed a language to write Unix · Direct hardware control"),
    ("1983", "C++ — Bjarne Stroustrup","C + objects · Organise code around real-world things"),
    ("1991", "Python — Guido van Rossum","Written over Christmas · Named after Monty Python · Built to be fun"),
    ("1995", "Java — Sun Microsystems","Write once, run anywhere · Runs on a JVM · Dominates enterprise"),
    ("1995", "JavaScript — Brendan Eich","Written in 10 days for Netscape · Now runs servers too (Node.js)"),
], C["s3"], sub="Session 3 — Programming Languages")

# 36. Why Python Won AI
make_content("Why Python Dominates AI and Data Science", [
    "📖  Readable — code reads almost like English. Researchers share it openly.",
    "📦  Libraries — NumPy (1995), Pandas (2008), PyTorch (2016). Built for Python first.",
    "🌐  Community — AI research standardised on Python. All papers have Python code.",
    "⚡  Fast enough — NumPy runs in C under the hood. Python is the interface, not the bottleneck.",
    "🎓  Guido's rule: 'Code is read more often than it is written.' Python optimises for reading.",
    "🏆  Result: learn Python and you can contribute to every AI library and research project",
], C["s3"], note="The network effect: Python won AI → all tutorials are Python → new researchers learn Python → all new libraries are Python. Unbreakable cycle.",
sub="Session 3 — Programming Languages")

# 37. Compiled vs Interpreted
make_comparison(
    "Compiled vs Interpreted",
    "⚙️ Compiled (C, C++, Rust)",
    ["Translated to machine code before running",
     "Fast at execution time",
     "Catch errors before the program runs",
     "Single target platform",
     "Used for: operating systems, games, embedded systems",
     "Edit → compile → run → test (slower cycle)"],
    "🐍 Interpreted (Python, JavaScript)",
    ["Translated line by line while running",
     "Slower, but flexible and instant",
     "Errors surface at runtime",
     "Runs anywhere with the interpreter",
     "Used for: data science, web, scripting, AI",
     "Edit → run → see result immediately (fast cycle)"],
    C["s3"], sub="Session 3 — Programming Languages"
)

# 38. Lab + Assignment
make_lab("Session 3", "Fibonacci — First Standalone Python File",
    ["Create a new file: fibonacci.py",
     "Write fibonacci_loop(n): using a for loop",
     "Write fibonacci_recursive(n): function calling itself",
     "Run both and compare the output",
     "Add the Golden Ratio section — print ratios of consecutive terms"],
    "fibonacci.py runs without errors.\nBoth functions produce the same output.\nGolden ratio converges to 1.618...",
    C["s3"],
    sports_note="Fibonacci squad numbers: 1,2,3,5,8,13,21,34,55 — find real players who've worn these")

# ── SESSION 4: PYTHON (slides 39–47) ────────────────────────────
print("  [Session 4] Slides 39–47")

# 39. Divider
make_chapter_break("4", "🐍", "Python",
    "Math with Superpowers",
    '"Python was named after Monty Python.\nIts inventor wanted programming to be fun.\nLet\'s test that claim."',
    C["s4"])

# 40. Python = Math You Already Know
make_two_col(
    "Python Is 12th-Grade Math — Just Different Notation",
    "📐 Math Notation",
    ["Variable:  x = 5",
     "Function:  f(x) = x²",
     "Set:  A = {1, 2, 3}",
     "Sequence: a₁, a₂, a₃, ...",
     "Summation: Σaᵢ for i in A",
     "Condition: if x > 0"],
    "🐍 Python Equivalent",
    ["x = 5",
     "def f(x): return x**2",
     "A = {1, 2, 3}",
     "a = [a1, a2, a3, ...]",
     "sum(a) or sum(a for a in A)",
     "if x > 0:"],
    C["s4"], left_col=PURPLE, right_col=C["s4"],
    sub="Session 4 — Python"
)

# 41. Core: Variables, Lists, Loops
make_code_slide("Variables, Lists, Loops — The Building Blocks",
    ["# Variables",
     "player = 'Vini Jr'   # string",
     "goals  = 24          # int",
     "xg     = 18.3        # float",
     "",
     "# Lists",
     "squad  = ['Vini', 'Bellingham', 'Mbappe']",
     "squad[0]        # → 'Vini'  (0-indexed!)",
     "len(squad)      # → 3",
     "",
     "# Loops",
     "for p in squad:",
     "    print(p)",
    ],
    "Variables store values.\nLists store ordered sequences of values.\nLoops repeat an action for each item.\n\n⚠️ Python is 0-indexed:\nfirst item = index 0\nsecond item = index 1\n\nThis trips everyone up at first — it's normal.",
    C["s4"], sub="Session 4 — Python")

# 42. Functions + Imports
make_code_slide("Functions and Imports — Reuse and Power",
    ["# Functions — exactly like f(x) in math",
     "def rate_season(goals, assists):",
     "    total = goals + assists",
     "    if total > 30: return 'Exceptional'",
     "    elif total > 20: return 'Good'",
     "    else: return 'Average'",
     "",
     "print(rate_season(24, 11))  # → 'Exceptional'",
     "",
     "# Imports — use other people's code",
     "import numpy as np",
     "import matplotlib.pyplot as plt",
     "x = np.linspace(0, 10, 100)",
     "plt.plot(x, np.sin(x))",
     "plt.show()",
    ],
    "def creates a reusable block — just like defining f(x) in math.\n\nImports let you use libraries:\n\nnumpy → math operations on arrays (like matrices)\n\nmatplotlib → graphs and charts\n\npandas → tables of data (like Excel in Python)\n\nThese three libraries power almost all data science.",
    C["s4"], sub="Session 4 — Python")

# 43. Lab 1: Math Waves
make_lab("Session 4 — Lab 1", "Math Waves — Plot Functions from 12th Grade",
    ["Open labs/lab1_math_waves.py",
     "Run Part 1 — Python as a calculator (verify familiar results)",
     "Run Part 2 — plot sin(x) and cos(x) on the same graph",
     "Understand every line: np.linspace, plt.plot, plt.show",
     "Run Part 3 — plot the quadratic x²-5x+6 with roots marked"],
    "Two matplotlib windows open:\n1. Sin and cos wave plot\n2. Quadratic with roots at x=2 and x=3 highlighted",
    C["s4"],
    sports_note="Challenge: plot the derivative of sin(x). It should look like cos(x) — verify visually")

# 44. Lab 2: Guessing Game
make_lab("Session 4 — Lab 2", "The Guessing Game — Binary Search in Action",
    ["Open labs/lab2_guessing_game.py",
     "Play the game: find the secret number, count your attempts",
     "Try to beat it in 7 guesses (it's possible — here's how)",
     "Watch binary_search_strategy() run — always guess the middle",
     "Sree explains: this is O(log n) — 1,000 numbers in 10 guesses"],
    "Play the game yourself first.\nWatch the binary search demo.\nUnderstand why log₂(1000) ≈ 10.",
    C["s4"],
    sports_note="Binary search is how NBA databases find any player instantly among millions of records")

# 45. Lab 3: NBA + Soccer Charts
make_lab("Session 4 — Lab 3", "Top Scorers Dashboard — NBA + Soccer",
    ["Open labs/lab3_top_scorers.py",
     "FIRST: change all player names and stats to your actual favourites",
     "Run it — two charts appear side by side",
     "Left chart: soccer goals + assists bar chart",
     "Right chart: basketball PPG vs APG scatter plot"],
    "Two charts saved as top_scorers.png\nAdvaith's own favourite players, his own data.\nSend it to friends — it looks professional.",
    C["s4"],
    sports_note="Use YOUR players — Mbappe, Luka, Vini, Giannis — whoever you actually watch")

# 46. Lab 4: Dice Simulator
make_lab("Session 4 — Lab 4", "Dice Simulator — Law of Large Numbers",
    ["Open labs/lab4_dice_simulator.py",
     "Run it — 4 histograms appear (10, 100, 1000, 10000 rolls)",
     "Observe: 10 rolls looks lumpy, 10,000 rolls looks flat",
     "Read the free throw simulation at the bottom",
     "Discussion: what does this tell us about 5-game slumps in sports?"],
    "4 histograms in a grid — showing convergence to flat distribution.\nFree throw simulation showing how sample size affects accuracy.",
    C["s4"],
    sports_note="Why a 5-game slump means nothing statistically. Why a full 82-game NBA season tells you everything.")

# 47. Session 4 Assignment
make_assignment("Session 4 — Python", [
    "Task 1: Pick TWO lab challenge prompts and complete them. Bring the code and output.",
    "Task 2: Build a match stats tracker — input 5 results, print wins/draws/losses, goals scored/conceded, goal difference, points.",
    "Task 3: Update lab3_top_scorers.py to include 10 players (5 soccer + 5 basketball). Add minutes played and compute goals per 90.",
    "Task 4: How many guesses does binary search need for 1 billion records? Why does Google Search return results in 0.5 seconds? Connect to Big-O.",
], C["s4"])

# ── SESSION 5: MATH FOR AI (slides 48–53) ───────────────────────
print("  [Session 5] Slides 48–53")

# 48. Divider
make_chapter_break("5", "🧠", "Math & Stats for AI",
    "The Engine Behind Every AI",
    '"Netflix, Spotify, ChatGPT — they all reduce to three things\nfrom your 12th-grade textbook: calculus, linear algebra, probability."',
    C["s5"])

# 49. Linear Algebra for AI
make_content("Vectors and Matrices — Neural Networks Are Matrix Math", [
    "📐  A vector = a list of numbers. Player stats: [goals, assists, xG, minutes] = [24, 11, 18.3, 2650]",
    "🧮  Dot product: [1,2,3]·[4,5,6] = 1×4 + 2×5 + 3×6 = 32. Core operation of every neural network.",
    "📊  Matrix = a grid of numbers. A neural network 'layer' = matrix × input vector + bias",
    "🔄  Matrix multiplication: transform one vector into another. Each layer reshapes the data.",
    "⚡  The entire forward pass of GPT-4 = billions of matrix multiplications happening in parallel on GPUs",
    "✅  Advaith already knows vectors from 12th-grade physics. He already knows matrix multiplication.",
], C["s5"], note="The whole secret of neural networks: output = matrix × input + bias. Repeated many times. That's it.",
sub="Session 5 — Math & Stats for AI")

# 50. Gradient Descent
make_content("Gradient Descent — How AI Learns", [
    "🎯  Learning = finding parameters that minimise a 'loss function' (how wrong we are)",
    "📉  Loss function measures the gap between prediction and truth",
    "📐  The derivative tells us the slope at any point — which direction is downhill",
    "⬇️  Move the parameters in the downhill direction → predictions improve",
    "🔁  Repeat thousands of times → the model 'rolls downhill' to the minimum",
    "🌍  ChatGPT's training: gradient descent in a space with 175 billion dimensions, for months",
], C["s5"], note="Sports analogy: gradient descent is a player watching game film, finding one specific weakness, correcting it, repeat. Each correction = one gradient step.",
sub="Session 5 — Math & Stats for AI")

# 51. Probability + Stats
make_two_col(
    "Probability & Statistics — How AI Decides",
    "📊 Key Stats Concepts",
    ["Mean — average outcome (expected goals per season)",
     "Variance — spread around the mean (consistency)",
     "Std deviation — √variance (volatility = risk)",
     "Correlation — how two variables move together",
     "Normal distribution — the bell curve, appears everywhere",
     "Bayes' Theorem — update probability when new evidence arrives"],
    "⚽🏀 Sports Applications",
    ["xG — expected goals is a probability: 0.0 → 1.0",
     "Free throw % — probability of making each shot",
     "'He's due a goal' — Bayes says NO. Each shot is independent.",
     "Injury prediction — correlation between workload and injury risk",
     "Player valuation — regression model on stats and performance",
     "Scouting — find players whose stats exceed their reputation (undervalued)"],
    C["s5"], sub="Session 5 — Math & Stats for AI"
)

# 52. Lab: xG Model
make_lab("Session 5", "The xG (Expected Goals) Model",
    ["Open labs/xg_model.py — read the shot data at the top",
     "Run Part 2: plot Goals vs Distance (scatter plot)",
     "Run Part 2: plot Distance + Angle coloured by outcome",
     "Identify the pattern — draw a curve through the data by hand",
     "Run Part 3: use estimate_xg() on famous shots",
     "Check: what does the model say about Mbappe vs Argentina?"],
    "Two scatter plots showing the xG pattern.\nFamous shot analysis printed.\nMbappe's goal: xG = 0.08 (8% chance — he scored it anyway).",
    C["s5"],
    sports_note="That 0.08 xG goal is why Mbappe is Mbappe. The model quantifies what the eye already knew.")

# 53. Assignment
make_assignment("Session 5 — Math & Stats for AI", [
    "Task 1: Extend estimate_xg() with two new features (e.g., header vs foot, counter-attack). Test on 5 famous shots. Did scores change?",
    "Task 2: Calculate dot products by hand: [3,4]·[1,2], [1,0,0]·[5,3,7], [2,3,1]·[2,3,1]. Verify with numpy.dot(). What is a vector dotted with itself?",
    "Task 3: Compute correlation between goals and xG for 10 players. Plot goals vs xG scatter. Add y=x line. Who overperforms xG most?",
    "Task 4: In 2 sentences, explain local minimum vs global minimum using a basketball analogy.",
], C["s5"])

# ── SESSION 6: AGENTIC AI (slides 54–60) ────────────────────────
print("  [Session 6] Slides 54–60")

# 54. Divider
make_chapter_break("6", "🤖", "Agentic AI",
    "The Thing That Changes Everything",
    '"In 2022, an AI passed the bar exam and aced the SAT.\nIn 2025, AI agents book meetings, write code, and run workflows.\nYou\'re building one today."',
    C["s6"])

# 55. The Arc: Rules → ML → Deep Learning
make_timeline("70 Years of AI — The Story So Far", [
    ("1950s", "Rule-Based AI",       "Programmers write every rule by hand · Works for chess, fails for reality"),
    ("1990s", "Statistical ML",      "Algorithms learn patterns from examples · Spam filters, credit scoring"),
    ("2012",  "Deep Learning",       "AlexNet wins ImageNet · Neural nets learn their own features · Game changer"),
    ("2017",  "'Attention Is All You Need'",
                                     "Transformer architecture · Parallel training · Scales to any size"),
    ("2022",  "ChatGPT",            "1M users in 5 days · 100M in 2 months · Fastest adoption in history"),
    ("2025",  "AI Agents",          "AI that reasons, uses tools, takes actions · The era Advaith enters"),
], C["s6"], sub="Session 6 — Agentic AI")

# 56. Transformers + LLMs
make_content("Transformers — Why Everything Changed in 2017", [
    "🧠  Attention mechanism — model learns which parts of input to focus on for each output token",
    "⚡  Parallel training — unlike previous RNNs, Transformers train all tokens simultaneously → can use GPUs",
    "📈  Scale — bigger model + more data + more compute → dramatically better results (scaling laws)",
    "🌍  GPT-3 (2020): 175 billion parameters · Trained on most of the internet",
    "💬  Emergent behaviour — nobody programmed GPT-3 to translate or write code · It learned from scale",
    "🏗️  Every major AI today — GPT, Claude, Gemini, Llama — is a Transformer",
], C["s6"], note="The 2017 paper 'Attention Is All You Need' by Google Brain has been cited 100,000+ times. It changed AI as completely as transistors changed hardware.",
sub="Session 6 — Agentic AI")

# 57. What Makes an Agent
make_two_col(
    "Chatbot vs Agent — A Fundamental Difference",
    "💬 Chatbot",
    ["Receives a question",
     "Generates an answer",
     "One turn → done",
     "Passive — waits for instructions",
     "No tools, no memory beyond the conversation",
     "Like texting someone a question"],
    "🤖 Agent",
    ["Receives a goal",
     "Plans the steps to achieve it",
     "Uses tools (search, code, APIs)",
     "Has memory across sessions",
     "Takes real actions in the world",
     "Like hiring a contractor to get something done"],
    C["s6"], left_col=MGREY, right_col=C["s6"],
    sub="Session 6 — Agentic AI"
)

# 58. The Claude API
make_code_slide("The Claude API — 8 Lines to a Working AI",
    ["import anthropic",
     "",
     "client = anthropic.Anthropic()",
     "",
     "response = client.messages.create(",
     '    model="claude-opus-4-7",',
     "    max_tokens=300,",
     '    system="You are a sharp sports analyst.",',
     "    messages=[{",
     '        "role": "user",',
     '        "content": "Analyse Vinícius Jr"',
     "    }]",
     ")",
     "print(response.content[0].text)",
    ],
    "model = the brain\nsystem = the personality\nmessages = the memory\nmax_tokens = maximum length\n\nChange system from 'sharp analyst' to 'harsh critic' to 'fanboy' — same model, completely different personality.\n\nThe system prompt is a dial. Experiment with it.",
    C["s6"], sub="Session 6 — Agentic AI")

# 59. Lab: Sports Analyst
make_lab("Session 6", "Build the Sports Analyst Agent",
    ["Open labs/sports_analyst.py — set ANTHROPIC_API_KEY first",
     "Run Part 1: one-shot analysis of any player",
     "Run Part 2: same player, three different personas — compare",
     "Read Part 3: understand conversation_history (this IS memory)",
     "Run the scripted multi-turn demo — watch it remember 'him'",
     "Challenge: add debate mode — agent argues Player A > Player B"],
    "Working sports analyst responding to real queries.\nPersona experiment: 3 different voices, same player.\nMulti-turn demo showing memory across turns.",
    C["s6"],
    sports_note="Ask it about Mbappe, Luka, Vini, Bellingham — then ask it to compare them to legends at the same age")

# 60. Assignment
make_assignment("Session 6 — Agentic AI", [
    "Task 1: Add debate_mode(player1, player2) — agent picks a side, gives 3 arguments. Test: Mbappe vs Haaland, Luka vs SGA.",
    "Task 2: Feed your Session 4 player stats into the agent. Ask: who is over/under-performing their xG? Is the AI's answer correct?",
    "Task 3: Try 3 very different system prompts for the same question. Which was most useful? Which was most entertaining?",
    "Task 4 (Think): GPT-4 passed the bar exam. Does that mean it 'understands' law? Write 4–6 sentences on your view.",
], C["s6"])

# ── SESSION 7: QUANT FINANCE (slides 61–67) ─────────────────────
print("  [Session 7] Slides 61–67")

# 62. Divider
make_chapter_break("7", "📈", "Quant Finance",
    "Hedge Funds Are Just Applied Math",
    '"In 1988, mathematicians left academia to start a hedge fund.\nNo traders. No finance degrees. Just math people.\nThey made $100 billion."',
    C["s7"])

# 62. Renaissance Technologies Story
make_content("Renaissance Technologies — The Greatest Trade Ever Made", [
    "🧮  Jim Simons: mathematician, Cold War codebreaker, MIT professor",
    "🏦  1988: founded Renaissance Technologies — hired physicists and mathematicians, no traders",
    "📊  Medallion Fund: 66% average annual return for 30 years (before fees)",
    "🌍  S&P 500 average: ~10% per year. Medallion: 66%. No other fund comes close.",
    "🔐  Medallion closed to outside investors in 1993 — too profitable to share",
    "🧩  Secret: treat financial markets like a physics problem. Find patterns others can't see.",
], C["s7"], note="The same skills: linear algebra, probability, Python, time series analysis — used in AI and in quant finance. The tools are identical. The data is different.",
sub="Session 7 — Quant Finance")

# 63. Returns + Volatility
make_code_slide("Returns and Volatility — The Two Core Measures",
    ["import yfinance as yf",
     "import numpy as np",
     "",
     "# Download 3 years of Man United data",
     "manu = yf.Ticker('MANU').history(period='3y')['Close']",
     "",
     "# Daily return = % change each day",
     "returns = manu.pct_change().dropna()",
     "",
     "# Annualise (252 trading days/year)",
     "annual_return = returns.mean() * 252",
     "annual_vol    = returns.std() * np.sqrt(252)",
     "",
     "print(f'Return:     {annual_return:.1%}')",
     "print(f'Volatility: {annual_vol:.1%}')",
    ],
    "Return = how much you made.\nVolatility = how much it swings.\n\nVolatility = standard deviation of returns.\n\nHigh volatility = risky investment.\n\nA consistent player scoring 20 goals every season has LOW volatility.\n\nA player scoring 35 one year and 8 the next has HIGH volatility.\n\nWould you sign the volatile player? Depends on the price.",
    C["s7"], sub="Session 7 — Quant Finance")

# 64. Portfolio Theory + Sharpe Ratio
make_content("Markowitz + Sharpe — Two Nobel Prize Ideas", [
    "📊  Portfolio theory (Markowitz, 1952): diversification reduces risk without reducing return",
    "🔗  If two assets are uncorrelated, holding both lowers total portfolio volatility",
    "➗  Sharpe Ratio = (Return − Risk-Free Rate) / Volatility",
    "📈  Sharpe > 1.0 = good · Sharpe > 2.0 = excellent · Renaissance Medallion = ~2.0",
    "🌍  S&P 500 Sharpe ≈ 0.5. Most hedge funds: 0.3–0.8. Renaissance: off the chart.",
    "⚽  Transfer budget as a portfolio: buy uncorrelated players (different positions, risk profiles)",
], C["s7"], note="The Sharpe Ratio appears in every investment discussion. Learn to interpret it: 'How much return am I getting per unit of risk I'm taking?'",
sub="Session 7 — Quant Finance")

# 65. Black-Scholes
make_content("Black-Scholes — The Equation That Won the Nobel Prize", [
    "📜  Fischer Black and Myron Scholes (1973): formula to price options on stocks",
    "❓  An option = the right (not obligation) to buy a stock at a fixed price in the future",
    "🧮  The formula uses: logarithms, exponentials, normal distribution, partial derivatives",
    "✅  Every single symbol in Black-Scholes is something Advaith already knows from 12th grade",
    "💰  This formula governs a $10 trillion global derivatives market",
    "🎓  Scholes won the Nobel Prize in Economics for it (Black had died by then)",
], C["s7"], note="Sree: write the Black-Scholes formula on a whiteboard. Point to each symbol. 'You know this. This. And this. All of it.' This is the moment.",
sub="Session 7 — Quant Finance")

# 66. Lab: Stock Analysis
make_lab("Session 7", "Man United Stock Analysis — Real Financial Data",
    ["Open labs/stock_analysis.py and run it",
     "Watch MANU, AAPL, NVDA download — 3 years of real data",
     "Observe the MANU chart: find the Nov 2022 Glazers-for-sale spike",
     "Read the Sharpe Ratio table — which has the best risk-adjusted return?",
     "Study the correlation matrix — which two assets are least correlated?"],
    "Three price history charts saved.\nSharpe Ratio comparison table printed.\nCorrelation matrix printed.",
    C["s7"],
    sports_note="Man United: same numbers, two answers. Good for Glazers (3× return). Bad for fans (debt killed transfers).")

# 67. Assignment
make_assignment("Session 7 — Quant Finance", [
    "Task 1: Add two new tickers (Nike NKE, EA Sports, MSG Sports). Calculate their Sharpe Ratio. Does adding them to a portfolio with MANU reduce risk?",
    "Task 2: Transfer Budget Model — 3 striker options with fees, goals, injury days. Compute mean goals, std deviation, goals/£M. Build a Sharpe-like ratio. Who do you sign?",
    "Task 3: Research: what is Renaissance Technologies' average annual return? Why did they close to outside investors in 1993?",
    "Task 4: How is building an xG model similar to a trading strategy? How could your Session 6 agent be used in quant trading? What risks would you worry about?",
], C["s7"])

# ── SESSION 8: CAPSTONE (slides 68–73) ──────────────────────────
print("  [Session 8] Slides 68–73")

# 68. Divider
make_chapter_break("8", "🏆", "Capstone Project",
    "Build Something Real",
    '"Eight sessions of learning.\nOne rule: every line of code you write, you can explain.\nThat\'s all."',
    C["s8"])

# 69. Track A: The Oracle
make_content("Track A — The Oracle: AI Stock Analysis Agent", [
    "📊  Input: a stock ticker. Output: an AI-generated investment thesis.",
    "📥  Download price data with yfinance",
    "🧮  Calculate: annual return, volatility, Sharpe Ratio",
    "🤖  Feed stats to Claude API → generates 3-paragraph plain-English analysis",
    "💡  Sessions used: 4 (Python), 6 (Claude API), 7 (finance metrics)",
    "🏁  Demo: type 'MANU' → full analysis appears in seconds",
], C["s8"], note="3-day build: Day 3 → get_stock_stats() working · Day 4 → Claude API connected · Day 5 → polish and error handling",
sub="Session 8 — Capstone")

# 70. Track B: NeuralNet from Zero
make_content("Track B — NeuralNet from Zero: Digit Recognition", [
    "🖊️  Train a neural network to recognise handwritten digits (0–9)",
    "🚫  Rules: NO PyTorch, NO TensorFlow, NO scikit-learn",
    "✅  Tools allowed: NumPy only — every operation is visible, explainable math",
    "🧠  Architecture: 784 inputs (28×28 pixels) → 64 hidden units → 10 outputs",
    "🎯  Target accuracy: >85% on test set",
    "💡  Sessions used: 4 (Python + NumPy), 5 (gradient descent, matrix math)",
], C["s8"], note="This is the hardest track. But if Advaith can explain how a neural network learns from scratch, he understands something most CS graduates can't.",
sub="Session 8 — Capstone")

# 71. Track C: Tutor Bot
make_content("Track C — The Tutor Bot: AI That Teaches Math", [
    "📚  User chooses a topic: derivatives, probability, linear algebra, etc.",
    "🤖  AI explains the concept at the right level (system prompt controls depth)",
    "📝  AI generates a practice problem tailored to the topic",
    "✅  User submits an answer → AI gives feedback and adjusts difficulty",
    "🔁  Loop continues — getting harder or easier based on performance",
    "💡  Sessions used: 4 (Python), 6 (Claude API + multi-turn memory)",
], C["s8"], note="Meta-moment: Advaith is building a system that does what Sree has been doing. The student builds the teacher.",
sub="Session 8 — Capstone")

# 72. Track D: Sports Oracle (Recommended)
make_content("Track D — The Sports Oracle  ⭐ Recommended", [
    "📊  Load real soccer + basketball player data (CSV provided)",
    "🧮  Compute: goal contributions, goals per 90, xG over/under-performance, efficiency",
    "📈  Plot: 3 charts — goals bar chart, goals vs assists scatter, xG comparison",
    "🤖  Claude API summarises the 3 most interesting findings in the data",
    "🏁  Demo: running program → charts open → AI commentary prints below",
    "💡  Sessions used: 2 (stats), 4 (Python + pandas), 5 (xG), 6 (Claude API)",
], C["s8"], note="Starter file: sports_dashboard_starter.py — skeleton with empty functions. Advaith fills them in using what he learned.",
sub="Session 8 — Capstone")

# 73. 7-Day Build + Demo Day
make_two_col(
    "How the Capstone Week Works",
    "🗓️ 7-Day Structure",
    ["Day 1–2: Scope and scaffold with Sree",
     "→ Pick track, define functions, create skeleton",
     "→ Advaith leaves with something that runs (even if empty)",
     "Day 3–5: Advaith builds independently",
     "→ Sree available to unblock — NOT to write code",
     "Day 6: Code review session",
     "→ Sree reviews like a senior engineer",
     "→ One improvement implemented together",
     "Day 7: Demo Day"],
    "🎤 Demo Day (5 minutes)",
    ["0:00–2:00  Live demo — show the full program running",
     "2:00–3:00  Explain one piece of code line by line",
     "3:00–4:00  What surprised you during building?",
     "4:00–5:00  What would you add if you had one more week?",
     "",
     "Sree asks 3 questions:",
     "→ Walk me through what happens when [input]",
     "→ What would break this with different data?",
     "→ Which session did you use most?"],
    C["s8"], sub="Session 8 — Capstone"
)

# ── SESSION 9: SOFTWARE ENGINEERING (slides 74–80) ───────────────
print("  [Session 9] Slides 74–80")

# 74. Divider
make_chapter_break("9", "🏗️", "Software Engineering Landscape",
    "How the Profession Evolved: 2010 → 2026",
    '"In 2010 you deployed by copying files to a server manually.\nIn 2026 AI writes your code and your infrastructure auto-heals.\nSame job title. Completely different world."',
    C["s9"])

# 75. Eras 2010–2014
make_two_col(
    "The Foundation Years — 2010 to 2014",
    "2010  Java + OOP + Web Fundamentals",
    ["Java: the dominant enterprise language — still runs most banking",
     "OOP: organise code around objects (a Player has goals and can shoot())",
     "HTML/CSS/JavaScript: what users see, style, and interact with",
     "Server-side: request → query database → return HTML",
     "The model: learn one language deeply. Mastery mattered."],
    "2014  SPA + TypeScript",
    ["SPA (Single-Page Application): page loads once, content swaps dynamically",
     "React / Angular / Vue: frameworks managing SPA complexity",
     "You saw SPAs in DevTools: those XHR calls are SPAs talking to servers",
     "TypeScript: JavaScript with type hints — catches bugs before running",
     "Frontend became a serious engineering specialisation",
     "Instagram, YouTube, Gmail — all SPAs"],
    C["s9"], sub="Session 9 — Software Engineering Landscape"
)

# 76. Eras 2015–2016
make_two_col(
    "Scaling Up — 2015 to 2016",
    "2015  Spring Boot + REST APIs",
    ["Spring Boot: a framework that wires up a backend with minimal config",
     "REST API: standard contract using HTTP verbs (GET, POST, PUT, DELETE)",
     "When your sports_analyst.py calls Claude, it makes a REST POST request",
     "Enterprise patterns: dependency injection, ORMs, middleware",
     "Backend = assembling components, not building from scratch"],
    "2016  Microservices + DevOps + Docker",
    ["Microservices: split one big app into many small independent services",
     "Netflix has ~700 microservices. Each team deploys independently.",
     "Docker: package app + dependencies into a container — runs anywhere",
     "CI/CD: push code → tests run automatically → deploy to production",
     "DevOps: developers own deployment. No more 'throw it over the wall.'",
     "Netflix deploys thousands of times per day"],
    C["s9"], sub="Session 9 — Software Engineering Landscape"
)

# 77. Eras 2017–2020
make_two_col(
    "At Scale — 2017 to 2020",
    "2017  Cloud Engineering",
    ["AWS / Azure / GCP: rent compute, storage, networking — pay per second",
     "No data centre, no hardware, no cooling — just an API",
     "Lambda: run code without managing a server (serverless)",
     "S3: store any file. RDS: managed database. CDN: files served globally.",
     "A 2-person startup can now serve millions on the same infra as Netflix",
     "Infrastructure as Code: describe your servers in code. Version-controlled."],
    "2020  Kubernetes + Observability",
    ["Kubernetes (K8s): declare 'I want 5 copies of this service running always'",
     "K8s handles: starting containers, restarting crashes, routing traffic, scaling",
     "Observability: when 700 microservices exist, how do you find what's broken?",
     "Logs: every service writes what it did",
     "Metrics: numbers over time — request rate, error rate, response time",
     "Tracing: follow one user request across all 12 services it touched"],
    C["s9"], sub="Session 9 — Software Engineering Landscape"
)

# 78. Eras 2023–2026
make_two_col(
    "The Present and Future — 2023 to 2026",
    "2023  Platform Engineering + Security",
    ["Every team rebuilding the same CI/CD pipeline = massive waste",
     "Platform team builds a 'golden path': opinionated, pre-built way to ship",
     "New developer joins → clone template → one command → full service",
     "Shift-left security: find vulnerabilities in the IDE, not in production",
     "Automated scanners check every pull request before it merges",
     "'Platform engineers build the tools that make everyone else go faster'"],
    "2026  AI-Native Engineering  ← You Are Here",
    ["Code copilots: AI writes first drafts, you review and direct",
     "Agentic workflows: AI takes multi-step actions autonomously",
     "LLM apps: AI is the core feature, not an add-on",
     "Engineers using AI tools write 2–5× more code per day",
     "Skills that matter more: judgment, system design, knowing when AI is wrong",
     "Advaith enters engineering when it is being reinvented. No habits to unlearn."],
    C["s9"], left_col=SLATE, right_col=RED,
    sub="Session 9 — Software Engineering Landscape"
)

# 79. The Full Stack Diagram
s = slide()
bg(s, DARK)
rect(s, 0, 0, Inches(0.15), H, C["s9"])
txt(s, "Every Request You Make Travels Through All of This",
    Inches(0.35), Inches(0.28), Inches(12.5), Inches(0.75),
    sz=26, bold=True, col=WHITE)
hline(s, Inches(1.15), C["s9"])

layers = [
    ("SPA / React",             "Era 3, 2014",  BLUE),
    ("REST API / API Gateway",  "Era 4, 2015",  PURPLE),
    ("Microservices",           "Era 5, 2016",  GREEN),
    ("Kubernetes Cluster",      "Era 7, 2020",  TEAL),
    ("Cloud (AWS / GCP)",       "Era 6, 2017",  AMBER),
    ("AI Layer (LLM Agent)",    "Era 9, 2026",  RED),
]
lh = Inches(0.72)
for i, (layer, era, col) in enumerate(layers):
    y = Inches(1.38) + i * lh
    rect(s, Inches(0.38), y, Inches(8.5), Inches(0.6), col)
    txt(s, layer, Inches(0.55), y + Inches(0.12),
        Inches(6.0), Inches(0.38), sz=14, bold=True, col=WHITE)
    txt(s, era, Inches(6.8), y + Inches(0.15),
        Inches(2.0), Inches(0.35), sz=10, col=WHITE, italic=True, align=PP_ALIGN.RIGHT)
    # Arrow
    if i < len(layers) - 1:
        rect(s, Inches(4.3), y + Inches(0.6), Inches(0.5), Inches(0.12), MGREY)
txt(s, "← When you opened Instagram this morning, your request passed through all of these layers.",
    Inches(9.2), Inches(1.5), Inches(4.0), Inches(5.0),
    sz=12, col=LGREY, italic=True)

# 80. What Success Looks Like — Final Slide
s = slide()
bg(s, DARK)
rect(s, 0, 0, W, Inches(0.12), RED)
rect(s, 0, H - Inches(0.12), W, Inches(0.12), TEAL)
rect(s, 0, 0, Inches(0.12), H, RED)
rect(s, W - Inches(0.12), 0, Inches(0.12), H, TEAL)

txt(s, "What Success Looks Like",
    Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8),
    sz=32, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
hline(s, Inches(1.45), RED, Inches(11), Inches(1.15))

quote = (
    "At the start of this curriculum, Advaith had never written a line of code.\n\n"
    "By Demo Day, he can sit down at a keyboard, build a Python program that\n"
    "downloads real data, computes statistics, plots charts, and calls an AI API —\n"
    "and explain every single line he wrote."
)
txt(s, quote, Inches(1.2), Inches(1.72), Inches(11.0), Inches(2.5),
    sz=17, col=WHITE, italic=True, align=PP_ALIGN.CENTER)

hline(s, Inches(4.38), TEAL, Inches(11), Inches(1.15))

befores = ["Zero lines of code", "Math = school exercises", "AI = magic black box", "Finance = news headlines"]
afters  = ["4 Python projects + AI agent", "Math = the language of computing", "AI = built one from an API", "Finance = Sharpe, Markowitz, real data"]

for i, (b, a) in enumerate(zip(befores, afters)):
    y = Inches(4.6) + i * Inches(0.52)
    txt(s, f"✗  {b}", Inches(0.6), y, Inches(5.8), Inches(0.46),
        sz=13, col=RGBColor(0xee,0x88,0x88), align=PP_ALIGN.LEFT)
    txt(s, "→", Inches(6.5), y, Inches(0.6), Inches(0.46),
        sz=16, col=MGREY, align=PP_ALIGN.CENTER)
    txt(s, f"✓  {a}", Inches(7.2), y, Inches(5.85), Inches(0.46),
        sz=13, col=RGBColor(0x88,0xee,0xaa), align=PP_ALIGN.LEFT)

txt(s, "⚽  🏀  Let's go, Advaith.",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5),
    sz=18, bold=True, col=GOLD, align=PP_ALIGN.CENTER)


# ── SESSION 10: HOW INSTAGRAM WORKS (slides 81–87) ──────────────
print("  [Session 10] Slides 81–87")

# 81. Chapter break
make_chapter_break("10", "📸", "How Instagram Works",
    "The Billion-Dollar App You Use Without Thinking",
    '"In 30 seconds of scrolling,\nhundreds of servers worked for you."',
    C["s10"])

# 82. The 30-second test
make_content("What Happens When You Scroll?",
    [
        "🔐  Authenticate — server checks your login token",
        "🧠  Rank — ML model scores thousands of candidate posts",
        "📦  Fetch — images loaded from CDN servers near you",
        "📊  Predict — algorithm guesses which posts you'll like",
        "⏱️  Measure — every pause, zoom, and skip is recorded",
        "💰  Auction — an ad slot is won by the highest bidder",
    ],
    C["s10"],
    note="All of this happens in under 200 milliseconds. Every time.",
    sub="Session 10 — How Instagram Works")

# 83. The photo post journey
make_timeline("The Journey of Your Photo (Post to Feed)", [
    ("Tap Post",    "Phone → API Gateway",    "HTTP POST with JPEG + caption + auth token"),
    ("Processing",  "Media Service",          "Creates 4 versions: 150px thumb, 640px, 1080px, 1920px Story"),
    ("Storage",     "Amazon S3",             "All 4 versions stored permanently  ·  100+ billion photos total"),
    ("CDN",         "Edge Distribution",     "URLs pushed to 200+ nodes globally — Mumbai serves Mumbai"),
    ("Database",    "PostgreSQL write",      "post ID, user ID, timestamp, caption, location"),
    ("Fan-out",     "Kafka queue",           "Notification task created for every follower in parallel"),
], C["s10"], sub="Session 10 — How Instagram Works")

# 84. The algorithm
make_two_col("The Feed Algorithm — How Posts Are Ranked",
    "What the model measures",
    [
        "Past likes — have you liked this account before?",
        "Watch time — how long do you pause on their posts?",
        "DM frequency — do you message them?",
        "Recency — when was this posted?",
        "Popularity — is it going viral globally?",
    ],
    "What the algorithm decides",
    [
        "Interest score: your history with this account (40%)",
        "Relationship score: closeness (DMs, tags, etc.) (30%)",
        "Recency score: freshness decay over time (20%)",
        "Popularity score: global engagement (10%)",
        "Every action you take is training data for this model",
    ],
    C["s10"],
    sub="Session 10 — How Instagram Works")

# 85. Posts vs Stories vs Reels
make_comparison("Three Content Types — Three Different Systems",
    "Posts  (permanent)",
    [
        "Stored forever in Amazon S3",
        "Indexed in PostgreSQL",
        "Shown in feed + profile grid",
        "Can get likes for years",
        "Low urgency — durability matters",
    ],
    "Stories / Reels  (ephemeral + viral)",
    [
        "Stories: 24-hour TTL, auto-deleted from DB + CDN",
        "Reels: multi-resolution video (360p / 720p / 1080p)",
        "Adaptive bitrate: phone switches quality on 4G vs WiFi",
        "Reels: separate recommendation engine (like TikTok)",
        "A Reel from a nobody can reach 1 million people",
    ],
    C["s10"],
    left_col=BLUE, right_col=C["s10"],
    sub="Session 10 — How Instagram Works")

# 86. Scale and microservices
make_two_col("How It Handles 500 Million Daily Users",
    "The scale (at peak)",
    [
        "1 million photo uploads per minute",
        "4.2 billion likes per day",
        "100 million Stories per day",
        "A single server handles ~1,000 req/sec",
        "Instagram uses thousands of servers",
        "Auto-scaling: more servers at 8 PM Friday",
    ],
    "The microservices",
    [
        "auth-service — checks your login token",
        "media-service — stores + serves photos",
        "feed-service — builds your ranked feed",
        "ml-ranking-service — scores posts",
        "ads-service — runs the ad auction",
        "notification-service — push alerts",
    ],
    C["s10"],
    sub="Session 10 — How Instagram Works")

# 87. The wow + lab assignment
make_lab("Session 10", "Build Instagram's Feed Algorithm in Python",
    steps=[
        "Create a user profile dict with past_likes, watch_time, dm_frequency",
        "Define 8 raw posts with from, age_minutes, global_likes",
        "Write rank_post() using weighted formula (interest 40%, relationship 30%...)",
        "Sort posts by score descending — that's your ranked feed",
        "Simulate CDN fetch with latency by nearest node",
        "Change weights and see how the feed order changes",
    ],
    output_desc="Ranked feed printed to terminal\nshowing why each post ranks where it does.\n\nLike an action → model updates,\nfeed re-ranks next scroll.",
    col=C["s10"],
    sports_note="Kevin Systrom built the first Instagram in 13 days using Python. It sold for $1 billion. Same language you've been learning all month.")

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════

out = "/Users/srmallip/projects/jumpstartm_c/curriculum_advaith.pptx"
prs.save(out)

import os
total = len(prs.slides)
size  = os.path.getsize(out) / 1024
print(f"\n✅  Saved: {out}")
print(f"   Slides: {total}")
print(f"   Size:   {size:.0f} KB")
if total != 87:
    print(f"\n⚠️  Expected 87 slides, got {total}")
